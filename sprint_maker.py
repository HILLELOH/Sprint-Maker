import datetime
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Function to create a slide with job data
def create_slide(prs, df):
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Remove all shapes from the slide (clean slate)
    for shape in slide.shapes:
        sp = shape
        slide.shapes._spTree.remove(sp._element)

    # Set background color to white and add a placeholder for pattern-like effect
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

    # Draw a dotted pattern-like effect manually across the entire slide
    pattern_color = RGBColor(255, 51, 204)
    dot_spacing = 0.5  # Moderate spacing between dots for density
    slide_width = 25.4  # Slide width in cm (standard 10 inches)
    slide_height = 19.05  # Slide height in cm (standard 7.5 inches)

    for x in range(0, int(slide_width / dot_spacing)):
        for y in range(0, int(slide_height / dot_spacing)):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,  # Use oval for a circular dot effect
                left=Cm(x * dot_spacing),
                top=Cm(y * dot_spacing),
                width=Cm(0.05),  # Smaller dot size
                height=Cm(0.05)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = pattern_color
            shape.line.color.rgb = pattern_color

    # Define position and size of text boxes
    top = Cm(1)
    height = Cm(1.03)

    # Define the initial right position for the first text box (right-to-left layout)
    right_position = Cm(24)  # Start from the right edge

    # Loop through each row in the DataFrame and create text boxes
    for i, row in df.iterrows():
        print(f"Row {i}: {row}")
        mission = row['משימה']
        name = row['שם']
        time = str(row['זמן']).strip()  # Ensure time is treated as a string

        # Create the Mission Index text box (furthest to the right)
        box_index = slide.shapes.add_textbox(right_position - Cm(1.02), top, Cm(1.02), height)
        text_frame_index = box_index.text_frame
        text_frame_index.text = str(i + 1)  # Mission index (1-based)
        text_frame_index.paragraphs[0].alignment = PP_ALIGN.RIGHT  # Align text to the right
        box_index.line.color.rgb = RGBColor(0, 0, 0)  # Black border
        box_index.line.width = Cm(0.05)

        right_position -= Cm(1.02 + 0.5)  # Adjust for spacing

        # Create the Mission text box
        box1 = slide.shapes.add_textbox(right_position - Cm(7.04), top, Cm(7.04), height)
        text_frame1 = box1.text_frame
        text_frame1.text = mission
        text_frame1.paragraphs[0].alignment = PP_ALIGN.RIGHT  # Align text to the right
        box1.line.color.rgb = RGBColor(0, 0, 0)  # Black border
        box1.line.width = Cm(0.05)

        right_position -= Cm(7.04 + 0.5)  # Adjust for spacing

        # Create the Time text box (next to Mission)
        box2 = slide.shapes.add_textbox(right_position - Cm(1.94), top, Cm(1.94), height)
        text_frame2 = box2.text_frame
        text_frame2.text = time
        text_frame2.paragraphs[0].alignment = PP_ALIGN.RIGHT
        box2.line.color.rgb = RGBColor(0, 0, 0)  # Black border
        box2.line.width = Cm(0.05)

        right_position -= Cm(1.94 + 0.5)  # Adjust for spacing

        # Create the Name text box (furthest to the left)
        box3 = slide.shapes.add_textbox(right_position - Cm(2.55), top, Cm(2.55), height)
        text_frame3 = box3.text_frame
        text_frame3.text = name
        text_frame3.paragraphs[0].alignment = PP_ALIGN.RIGHT
        box3.line.color.rgb = RGBColor(0, 0, 0)  # Black border
        box3.line.width = Cm(0.05)

        # Adjust the top position for the next row
        top += Cm(1.53)  # Add vertical space for the next row

        # Reset right_position for the next row
        right_position = Cm(24)


# Function to read CSV and generate PowerPoint
def create_ppt_from_csv(csv_file, ppt_file):
    prs = Presentation()

    # Read CSV file into DataFrame
    df = pd.read_csv(csv_file)
    df.columns = df.columns.str.strip()  # Clean up any spaces in the column names

    # Print cleaned column names and row data for debugging
    print(f"Columns in CSV (cleaned): {df.columns}")

    # Create the slide with the job data
    create_slide(prs, df)

    # Save the PowerPoint presentation
    prs.save(ppt_file)
    print(f"Presentation saved to {ppt_file}")

# Define file paths for the CSV and PowerPoint presentation
current_dir = os.path.abspath(os.getcwd())
print(current_dir)
csv_directory = f'{current_dir}\csv files'
ppt_directory = f'{current_dir}\sprints'

# Ensure the directories exist
os.makedirs(csv_directory, exist_ok=True)
os.makedirs(ppt_directory, exist_ok=True)

# Specify CSV file and generate unique PowerPoint file name with a timestamp
csv_file = os.path.join(csv_directory, 'jobs.csv')
ppt_file = os.path.join(ppt_directory, f"presentation_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")

# Create the PowerPoint presentation from the CSV data
create_ppt_from_csv(csv_file, ppt_file)