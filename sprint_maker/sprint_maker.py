import datetime
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches

# Function to create slide with jobs data
def create_slide(prs, df):
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Remove all shapes from the slide (clean slate)
    for shape in slide.shapes:
        sp = shape
        slide.shapes._spTree.remove(sp._element)
    
    # Define position and size of text boxes
    top = Inches(1)
    width = Inches(2.5)
    height = Inches(0.5)
    
    # Define the initial left position for the first text box (Mission)
    left_position = Inches(0.5)
    
    # Loop through each row in the DataFrame and create text boxes
    for i, row in df.iterrows():
        print(f"Row {i}: {row}")
        mission = row['משימה']
        name = row['שם']
        time = str(row['זמן']).strip()  # Ensure time is treated as a string
        
        # Create the Mission text box
        box1 = slide.shapes.add_textbox(left_position, top, width, height)
        text_frame1 = box1.text_frame
        text_frame1.text = f"Mission: {mission}"

        # Create the Name text box (next to Mission)
        left_position += width + Inches(0.5)  # Move to the right
        box2 = slide.shapes.add_textbox(left_position, top, width, height)
        text_frame2 = box2.text_frame
        text_frame2.text = f"Name: {name}"

        # Create the Time text box (next to Name)
        left_position += width + Inches(0.5)  # Move to the right again
        box3 = slide.shapes.add_textbox(left_position, top, width, height)
        text_frame3 = box3.text_frame
        text_frame3.text = f"Time: {time} hours"
        
        # Adjust the top position for the next row (0.5 inches below the current row)
        top += height + Inches(0.5)  # Add 0.5 inches of vertical space for the next row

        # Reset left_position for the next row (to start from the left again)
        left_position = Inches(0.5)

# Function to read CSV and generate PowerPoint
def create_ppt_from_csv(csv_file, ppt_file):
    prs = Presentation()
    
    # Read CSV file into DataFrame
    df = pd.read_csv(csv_file)
    df.columns = df.columns.str.strip()  # Clean up any spaces in the column names
    
    # Print cleaned column names and row data for debugging
    print(f"Columns in CSV (cleaned): {df.columns}")
    
    # Create the slide with the job data
    create_slide(prs, df)
    
    # Save the PowerPoint presentation
    prs.save(ppt_file)
    print(f"Presentation saved to {ppt_file}")

# Define file paths for the CSV and PowerPoint presentation
csv_directory = r'C:\\projects\\sprint_maker\\csv files'
ppt_directory = r'C:\\projects\\sprint_maker\\sprints'

# Ensure the directories exist
os.makedirs(csv_directory, exist_ok=True)
os.makedirs(ppt_directory, exist_ok=True)

# Specify CSV file and generate unique PowerPoint file name with a timestamp
csv_file = os.path.join(csv_directory, 'jobs.csv')
ppt_file = os.path.join(ppt_directory, f"presentation_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")

# Create the PowerPoint presentation from the CSV data
create_ppt_from_csv(csv_file, ppt_file)
